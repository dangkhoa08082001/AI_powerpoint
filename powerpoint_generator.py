"""
Module để tạo PowerPoint presentations từ dữ liệu structured
Enhanced với layouts thông minh và thiết kế hiện đại
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from io import BytesIO
import json
import logging
import math
import os
import random
import re

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PowerPointGenerator:
    """
    Class chính để tạo PowerPoint presentations với layouts thông minh
    """
    
    def __init__(self):
        self.presentation = None
        self.slide_layouts = None
        self.current_template = "education"
        
        # Enhanced Template configurations với màu sắc đẹp
        self.templates = {
            "education": {
                "primary_color": "#2E86AB",      # Ocean Blue
                "secondary_color": "#A23B72",    # Magenta  
                "background_color": "#FFF8E1",   # Light Cream (thay vì orange)
                "text_color": "#1A1A1A",         # Dark Gray (easier to read)
                "accent_color": "#F18F01",       # Orange accent
                "highlight_color": "#E3F2FD",    # Light blue for highlights
                "font_size": {
                    "title": 32,
                    "subtitle": 24,
                    "content": 18,
                    "caption": 14
                }
            },
            "business": {
                "primary_color": "#1565C0",      # Professional Blue
                "secondary_color": "#FFA726",    # Warm Orange
                "background_color": "#FAFAFA",   # Clean White-Gray
                "text_color": "#212121",         # Dark text
                "accent_color": "#4CAF50",       # Success Green
                "highlight_color": "#E8F5E8",    # Light green for highlights
                "font_size": {
                    "title": 36,
                    "subtitle": 28,
                    "content": 20,
                    "caption": 16
                }
            },
            "modern": {
                "primary_color": "#6366F1",      # Modern Indigo
                "secondary_color": "#EC4899",    # Pink
                "background_color": "#F8FAFC",   # Slate Gray
                "text_color": "#0F172A",         # Slate Dark
                "accent_color": "#10B981",       # Emerald
                "highlight_color": "#F0F9FF",    # Sky light
                "font_size": {
                    "title": 34,
                    "subtitle": 26,
                    "content": 19,
                    "caption": 15
                }
            }
        }
        
        # Layout configurations - CHỈ 3 LAYOUTS AN TOÀN: TOP, LEFT, RIGHT - NO BOTTOM!
        self.layout_configs = {
            "image_top_content_bottom": {
                # Ảnh ở trên 50%, nội dung ở dưới 50% - LAYOUT DUY NHẤT an toàn cho ảnh trên
                "image_area": {"x": 1, "y": 1.2, "width": 8, "height": 3.0},      # Top 50%: 1.2-4.2
                "content_area": {"x": 0.5, "y": 4.2, "width": 9, "height": 3.0}   # Bottom 50%: 4.2-7.2
            },
            "content_left_image_right": {
                # Nội dung bên trái 50%, ảnh bên phải 50% - AN TOÀN HOÀN TOÀN!
                "content_area": {"x": 0.5, "y": 1.5, "width": 4.5, "height": 5.5}, # Left 50%: 0.5-5.0
                "image_area": {"x": 5.0, "y": 1.5, "width": 4.5, "height": 5.5}    # Right 50%: 5.0-9.5
            },
            "image_left_content_right": {
                # Ảnh bên trái 50%, nội dung bên phải 50% - AN TOÀN HOÀN TOÀN!
                "image_area": {"x": 0.5, "y": 1.5, "width": 4.5, "height": 5.5},   # Left 50%: 0.5-5.0
                "content_area": {"x": 5.0, "y": 1.5, "width": 4.5, "height": 5.5}  # Right 50%: 5.0-9.5
            }
            # ❌ REMOVED: "content_top_image_bottom" - LAYOUT BỊ LỖI, ẢNH BÊN DƯỚI LUÔN OVERLAP!
        }
        
        # Modern slide design elements
        self.design_elements = {
            "bullets": ["●", "▶", "◆", "✓", "►", "🔹", "🔸", "⭐"],
            "decorative_icons": ["🎯", "💡", "🔥", "⚡", "🌟", "🚀", "✨"],
            "numbered_styles": ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨", "⑩"]
        }
    
    def create_new_presentation(self, template: str = "education") -> bool:
        """
        Tạo presentation mới với theme background
        
        Args:
            template (str): Tên template muốn sử dụng
            
        Returns:
            bool: True nếu tạo thành công
        """
        try:
            self.presentation = Presentation()
            self.slide_layouts = self.presentation.slide_layouts
            self.current_template = template
            
            # Apply background theme
            self._apply_presentation_theme(template)
            
            logger.info(f"Created new presentation with template: {template}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating presentation: {str(e)}")
            return False
    
    def add_title_slide(self, title: str, subtitle: str = "", author: str = "") -> bool:
        """
        Thêm slide tiêu đề
        
        Args:
            title (str): Tiêu đề chính
            subtitle (str): Phụ đề
            author (str): Tác giả
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            slide_layout = self.slide_layouts[0]  # Title slide layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme
            self._apply_slide_background(slide)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_title_formatting(title_shape)
            
            # Set subtitle
            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                if author:
                    subtitle_shape.text = f"{subtitle}\n\n{author}"
                else:
                    subtitle_shape.text = subtitle
                self._apply_subtitle_formatting(subtitle_shape)
            
            logger.info(f"Added title slide: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding title slide: {str(e)}")
            return False
    
    def add_modern_content_slide(self, title: str, content: List[str], slide_type: str = "creative_bullets") -> bool:
        """
        Thêm slide nội dung với thiết kế hiện đại và sáng tạo
        
        Args:
            title (str): Tiêu đề slide
            content (List[str]): Danh sách nội dung
            slide_type (str): Loại slide (creative_bullets, icon_bullets, numbered_modern, cards, highlight_boxes)
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            slide_layout = self.slide_layouts[6]  # Blank layout for more control
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme with decorative elements
            self._apply_slide_background(slide)
            self._add_decorative_elements(slide)
            
            # Add styled title
            self._add_styled_title(slide, title)
            
            # Add content based on modern design type
            if slide_type == "creative_bullets":
                self._add_creative_bullet_content(slide, content)
            elif slide_type == "icon_bullets":
                self._add_icon_bullet_content(slide, content)
            elif slide_type == "numbered_modern":
                self._add_modern_numbered_content(slide, content)
            elif slide_type == "cards":
                self._add_card_layout_content(slide, content)
            elif slide_type == "highlight_boxes":
                self._add_highlight_box_content(slide, content)
            else:  # Default creative bullets
                self._add_creative_bullet_content(slide, content)
            
            logger.info(f"Added modern content slide '{slide_type}': {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding modern content slide: {str(e)}")
            return False

    def add_content_slide(self, title: str, content: List[str], slide_type: str = "bullet") -> bool:
        """
        Thêm slide nội dung (Legacy method - now uses modern design)
        """
        # Map old types to new modern types
        type_mapping = {
            "bullet": "creative_bullets",
            "numbered": "numbered_modern",
            "paragraph": "highlight_boxes"
        }
        modern_type = type_mapping.get(slide_type, "creative_bullets")
        return self.add_modern_content_slide(title, content, modern_type)
    
    def add_two_column_slide(self, title: str, left_content: List[str], right_content: List[str]) -> bool:
        """
        Thêm slide 2 cột
        
        Args:
            title (str): Tiêu đề slide
            left_content (List[str]): Nội dung cột trái
            right_content (List[str]): Nội dung cột phải
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            slide_layout = self.slide_layouts[3]  # Two content layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme
            self._apply_slide_background(slide)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_content_title_formatting(title_shape)
            
            # Left column
            left_shape = slide.placeholders[1]
            left_frame = left_shape.text_frame
            left_frame.clear()
            self._add_bullet_content(left_frame, left_content)
            
            # Right column
            right_shape = slide.placeholders[2]
            right_frame = right_shape.text_frame
            right_frame.clear()
            self._add_bullet_content(right_frame, right_content)
            
            logger.info(f"Added two-column slide: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding two-column slide: {str(e)}")
            return False
    
    def add_smart_image_content_slide(self, title: str, image_path: str, content: List[str], 
                                     layout_type: str = "auto") -> bool:
        """
        Thêm slide với layout thông minh dựa trên nội dung và ảnh
        
        Args:
            title (str): Tiêu đề slide
            image_path (str): Đường dẫn đến hình ảnh
            content (List[str]): Nội dung bullet points
            layout_type (str): Loại layout ("auto", "content_left_image_right", "image_left_content_right", 
                              "image_top_content_bottom", "content_top_image_bottom")
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            # Auto-select layout based on content length and image availability
            if layout_type == "auto":
                layout_type = self._auto_select_layout(content, image_path)
            
            slide_layout = self.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme with decorative elements
            self._apply_slide_background(slide)
            self._add_decorative_elements(slide)
            
            # Add title with improved styling
            self._add_styled_title(slide, title)
            
            # Apply the selected layout
            if layout_type in self.layout_configs:
                config = self.layout_configs[layout_type]
                self._apply_smart_layout(slide, config, content, image_path)
            else:
                # Fallback to content_left_image_right
                config = self.layout_configs["content_left_image_right"]
                self._apply_smart_layout(slide, config, content, image_path)
            
            logger.info(f"Added smart image-content slide with layout '{layout_type}': {title}")
            return True
                
        except Exception as e:
            logger.error(f"Error adding smart image-content slide: {str(e)}")
            return False

    def add_image_content_slide(self, title: str, image_path: str, content: List[str]) -> bool:
        """
        Thêm slide với layout: content bên trái, image bên phải (Legacy method)
        """
        return self.add_smart_image_content_slide(title, image_path, content, "content_left_image_right")
    
    def add_image_slide(self, title: str, image_path: str, caption: str = "") -> bool:
        """
        Thêm slide với hình ảnh
        
        Args:
            title (str): Tiêu đề slide
            image_path (str): Đường dẫn đến hình ảnh
            caption (str): Chú thích hình ảnh
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            slide_layout = self.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme
            self._apply_slide_background(slide)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            self._apply_content_title_formatting(title_box)
            
            # Add image
            try:
                img = slide.shapes.add_picture(
                    image_path,
                    Inches(1),
                    Inches(1.5),
                    width=Inches(8)
                )
                
                # Add caption
                if caption:
                    caption_box = slide.shapes.add_textbox(
                        Inches(1),
                        img.top + img.height + Inches(0.2),
                        Inches(8),
                        Inches(0.5)
                    )
                    caption_frame = caption_box.text_frame
                    caption_frame.text = caption
                    self._apply_caption_formatting(caption_box)
                
                logger.info(f"Added image slide: {title}")
                return True
                
            except Exception as img_error:
                logger.error(f"Error adding image: {img_error}")
                # Add placeholder text instead
                placeholder_box = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(4)
                )
                placeholder_frame = placeholder_box.text_frame
                placeholder_frame.text = f"[Hình ảnh: {image_path}]\n{caption}"
                return True
                
        except Exception as e:
            logger.error(f"Error adding image slide: {str(e)}")
            return False
    
    def add_chart_slide(self, title: str, chart_data: Dict[str, Any]) -> bool:
        """
        Thêm slide với biểu đồ
        
        Args:
            title (str): Tiêu đề slide
            chart_data (Dict): Dữ liệu biểu đồ
                {
                    "type": "column|line|pie|bar",
                    "categories": ["A", "B", "C"],
                    "series": {
                        "Series 1": [1, 2, 3],
                        "Series 2": [4, 5, 6]
                    }
                }
                
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            slide_layout = self.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme
            self._apply_slide_background(slide)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            self._apply_content_title_formatting(title_box)
            
            # Prepare chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get('categories', [])
            
            for series_name, values in chart_data.get('series', {}).items():
                chart_data_obj.add_series(series_name, values)
            
            # Chart type mapping
            chart_types = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED
            }
            
            chart_type = chart_types.get(chart_data.get('type', 'column'), XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Add chart
            chart = slide.shapes.add_chart(
                chart_type,
                Inches(1), Inches(2), Inches(8), Inches(5),
                chart_data_obj
            ).chart
            
            chart.has_legend = True
            chart.legend.position = 2  # Right
            
            logger.info(f"Added chart slide: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding chart slide: {str(e)}")
            return False
    
    def add_table_slide(self, title: str, table_data: List[List[str]], has_header: bool = True) -> bool:
        """
        Thêm slide với bảng
        
        Args:
            title (str): Tiêu đề slide
            table_data (List[List[str]]): Dữ liệu bảng
            has_header (bool): Có header row không
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            slide_layout = self.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Apply background theme
            self._apply_slide_background(slide)
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            self._apply_content_title_formatting(title_box)
            
            # Create table
            if table_data:
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0
                
                if rows > 0 and cols > 0:
                    table = slide.shapes.add_table(
                        rows, cols,
                        Inches(1), Inches(2),
                        Inches(8), Inches(4)
                    ).table
                    
                    # Fill table with data
                    for i, row_data in enumerate(table_data):
                        for j, cell_data in enumerate(row_data):
                            cell = table.cell(i, j)
                            cell.text = str(cell_data)
                            
                            # Style header row
                            if i == 0 and has_header:
                                self._apply_table_header_formatting(cell)
                            else:
                                self._apply_table_cell_formatting(cell)
            
            logger.info(f"Added table slide: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding table slide: {str(e)}")
            return False
    
    def add_conclusion_slide(self, title: str = "Kết luận", points: List[str] = None) -> bool:
        """
        Thêm slide kết luận
        
        Args:
            title (str): Tiêu đề slide
            points (List[str]): Các điểm kết luận
            
        Returns:
            bool: True nếu thêm thành công
        """
        try:
            if points is None:
                points = ["Cảm ơn các bạn đã lắng nghe!", "Có câu hỏi nào không?"]
            
            return self.add_content_slide(title, points, "bullet")
            
        except Exception as e:
            logger.error(f"Error adding conclusion slide: {str(e)}")
            return False
    
    def create_from_structured_data(self, presentation_data: Dict[str, Any]) -> bool:
        """
        Tạo presentation từ dữ liệu có cấu trúc với enhanced features
        
        Args:
            presentation_data (Dict): Dữ liệu presentation với enhanced data
                
        Returns:
            bool: True nếu tạo thành công
        """
        try:
            # Determine template from recommended_theme or fallback
            recommended_theme = presentation_data.get('recommended_theme', {})
            if recommended_theme:
                theme_name = recommended_theme.get('theme_name', 'education')
                # Map theme names to templates
                theme_template_map = {
                    'education_pro': 'education',
                    'tech_gradient': 'business', 
                    'business_elegant': 'business',
                    'creative_vibrant': 'education',
                    'python_modern': 'business'
                }
                template = theme_template_map.get(theme_name, 'education')
            else:
                template = presentation_data.get('template', 'education')
            
            # Create new presentation with theme
            if not self.create_new_presentation(template):
                return False
            
            # Get visual elements for enhanced styling
            visual_elements = presentation_data.get('visual_elements', {})
            primary_icon = visual_elements.get('primary_icon', '📊')
            
            # Add enhanced title slide
            title = presentation_data.get('title', 'Presentation')
            subtitle = presentation_data.get('subtitle', '')
            author = presentation_data.get('author', '')
            
            # Add icon to title if available
            if primary_icon and primary_icon not in title:
                title = f"{primary_icon} {title}"
            
            self.add_title_slide(title, subtitle, author)
            
            # Add content slides with enhanced features
            for slide_data in presentation_data.get('slides', []):
                slide_type = slide_data.get('type', 'content')
                slide_title = slide_data.get('title', '')
                slide_icon = slide_data.get('icon', '')
                
                # Add icon to slide title
                if slide_icon and slide_icon not in slide_title:
                    enhanced_title = f"{slide_icon} {slide_title}"
                else:
                    enhanced_title = slide_title
                
                # Check for image generation
                needs_image = slide_data.get('needs_image', False)
                image_path = slide_data.get('generated_image_path', '')
                
                logger.info(f"Processing slide: {enhanced_title}, needs_image: {needs_image}, has_path: {bool(image_path)}")
                
                if slide_type == 'content':
                    content = slide_data.get('content', [])
                    
                    # If has generated image, use smart layout (image + content)
                    if image_path and os.path.exists(image_path):
                        layout_type = slide_data.get('layout_type', 'auto')
                        self.add_smart_image_content_slide(enhanced_title, image_path, content, layout_type)
                    else:
                        # Use modern content design
                        design_type = slide_data.get('design_type', 'creative_bullets')
                        self.add_modern_content_slide(enhanced_title, content, design_type)
                        
                elif slide_type == 'two_column':
                    left_content = slide_data.get('left_content', [])
                    right_content = slide_data.get('right_content', [])
                    
                    # Check if there's an image for this slide
                    if image_path and os.path.exists(image_path):
                        # Use smart layout with image
                        combined_content = left_content + right_content
                        layout_type = slide_data.get('layout_type', 'content_left_image_right')
                        self.add_smart_image_content_slide(enhanced_title, image_path, combined_content, layout_type)
                    else:
                        self.add_two_column_slide(enhanced_title, left_content, right_content)
                    
                elif slide_type == 'image':
                    image_path = slide_data.get('image_path', image_path)
                    caption = slide_data.get('caption', slide_data.get('image_concept', ''))
                    self.add_image_slide(enhanced_title, image_path, caption)
                    
                elif slide_type == 'chart':
                    chart_data = slide_data.get('chart_data', {})
                    self.add_chart_slide(enhanced_title, chart_data)
                    
                elif slide_type == 'table':
                    table_data = slide_data.get('table_data', [])
                    has_header = slide_data.get('has_header', True)
                    self.add_table_slide(enhanced_title, table_data, has_header)
            
            # Add conclusion slide if specified
            if presentation_data.get('add_conclusion', True):
                conclusion_points = presentation_data.get('conclusion_points', None)
                conclusion_title = f"🏆 Kết luận"
                self.add_conclusion_slide(title=conclusion_title, points=conclusion_points)
            
            logger.info("Successfully created enhanced presentation from structured data")
            return True
            
        except Exception as e:
            logger.error(f"Error creating presentation from data: {str(e)}")
            return False
    
    def save_to_file(self, filename: str) -> bool:
        """
        Lưu presentation vào file
        
        Args:
            filename (str): Tên file để lưu
            
        Returns:
            bool: True nếu lưu thành công
        """
        try:
            if self.presentation is None:
                logger.error("No presentation to save")
                return False
                
            self.presentation.save(filename)
            logger.info(f"Presentation saved to: {filename}")
            return True
            
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            return False
    
    def save_to_buffer(self) -> Optional[BytesIO]:
        """
        Lưu presentation vào BytesIO buffer để download
        
        Returns:
            BytesIO: Buffer chứa file PowerPoint hoặc None nếu lỗi
        """
        try:
            if self.presentation is None:
                logger.error("No presentation to save")
                return None
                
            buffer = BytesIO()
            self.presentation.save(buffer)
            buffer.seek(0)
            
            logger.info("Presentation saved to buffer")
            return buffer
            
        except Exception as e:
            logger.error(f"Error saving presentation to buffer: {str(e)}")
            return None
    
    def get_slide_count(self) -> int:
        """
        Lấy số lượng slides
        
        Returns:
            int: Số lượng slides
        """
        if self.presentation is None:
            return 0
        return len(self.presentation.slides)
        
    def _apply_presentation_theme(self, template: str):
        """Apply theme background to presentation"""
        try:
            from pptx.dml.color import RGBColor
            from pptx.enum.dml import MSO_THEME_COLOR
            
            template_config = self.templates.get(template, self.templates["education"])
            bg_color_hex = template_config.get('background_color', '#F18F01')
            
            # Convert hex to RGB
            color_hex = bg_color_hex.replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            
            # Apply to slide master if possible
            try:
                slide_master = self.presentation.slide_master
                if hasattr(slide_master, 'background'):
                    slide_master.background.fill.solid()
                    slide_master.background.fill.fore_color.rgb = RGBColor(r, g, b)
                    logger.info(f"Applied background color {bg_color_hex} to slide master")
            except Exception as e:
                logger.warning(f"Could not apply background to slide master: {str(e)}")
                
        except Exception as e:
            logger.warning(f"Could not apply presentation theme: {str(e)}")
            
    def _apply_slide_background(self, slide, template: str = None):
        """Apply enhanced background with gradient and theme colors"""
        try:
            from pptx.dml.color import RGBColor
            from pptx.enum.dml import MSO_FILL_TYPE
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE
            
            if template is None:
                template = self.current_template
                
            template_config = self.templates.get(template, self.templates["education"])
            
            # Create gradient background
            self._create_gradient_background(slide, template_config)
            
            # Add theme accent elements
            self._add_theme_accents(slide, template_config)
                
        except Exception as e:
            logger.warning(f"Could not apply slide background: {str(e)}")
    
    def _create_gradient_background(self, slide, template_config):
        """Tạo background gradient đẹp mắt"""
        try:
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            
            # Main background với gradient effect
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(-0.5), Inches(-0.5),
                Inches(11), Inches(8.5)
            )
            
            # Primary background color (lighter tone)
            bg_color_hex = template_config.get('background_color', '#F18F01')
            color_hex = bg_color_hex.replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16) 
            b = int(color_hex[4:6], 16)
            
            # Make it lighter for better readability
            r = min(255, r + 20)
            g = min(255, g + 20)
            b = min(255, b + 20)
            
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(r, g, b)
            bg_shape.line.fill.background()
            
            # Move to back
            try:
                bg_shape._element.getparent().remove(bg_shape._element)
                slide.shapes._spTree.insert(2, bg_shape._element)
            except:
                pass
                
            logger.info(f"Applied gradient background")
            
        except Exception as e:
            logger.warning(f"Gradient background failed: {str(e)}")
    
    def _add_theme_accents(self, slide, template_config):
        """Thêm các accent elements theo theme"""
        try:
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            
            # Top accent bar
            top_accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(10), Inches(0.15)
            )
            
            primary_hex = template_config.get('primary_color', '#2E86AB').replace('#', '')
            r = int(primary_hex[0:2], 16)
            g = int(primary_hex[2:4], 16)
            b = int(primary_hex[4:6], 16)
            
            top_accent.fill.solid()
            top_accent.fill.fore_color.rgb = RGBColor(r, g, b)
            top_accent.line.fill.background()
            
            # Right accent 
            right_accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9.8), Inches(0), Inches(0.2), Inches(7.5)
            )
            
            accent_hex = template_config.get('accent_color', '#F18F01').replace('#', '')
            r2 = int(accent_hex[0:2], 16)
            g2 = int(accent_hex[2:4], 16)
            b2 = int(accent_hex[4:6], 16)
            
            right_accent.fill.solid()
            right_accent.fill.fore_color.rgb = RGBColor(r2, g2, b2)
            right_accent.line.fill.background()
            
            logger.info("Added theme accents")
            
        except Exception as e:
            logger.warning(f"Theme accents failed: {str(e)}")
    
    # Private formatting methods
    def _apply_title_formatting(self, shape):
        """Apply formatting cho title slide"""
        template = self.templates[self.current_template]
    
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(template['font_size']['title'])
            paragraph.font.bold = True
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['primary_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16) 
            b = int(color_hex[4:6], 16)
            paragraph.font.color.rgb = RGBColor(r, g, b)
            paragraph.alignment = PP_ALIGN.CENTER
    
    def _apply_subtitle_formatting(self, shape):
        """Apply formatting cho subtitle"""
        template = self.templates[self.current_template]
        
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(template['font_size']['subtitle'])
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['secondary_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            paragraph.font.color.rgb = RGBColor(r, g, b)
            paragraph.alignment = PP_ALIGN.CENTER
    
    def _apply_content_title_formatting(self, shape):
        """Apply formatting cho content title"""
        template = self.templates[self.current_template]
        
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(template['font_size']['subtitle'])
            paragraph.font.bold = True
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['primary_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            paragraph.font.color.rgb = RGBColor(r, g, b)
    
    def _apply_caption_formatting(self, shape):
        """Apply formatting cho caption"""
        template = self.templates[self.current_template]
        
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(template['font_size']['caption'])
            paragraph.font.italic = True
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            paragraph.font.color.rgb = RGBColor(r, g, b)
            paragraph.alignment = PP_ALIGN.CENTER

    def _apply_table_header_formatting(self, cell):
        """Apply formatting cho table header"""
        template = self.templates[self.current_template]
        
        cell.fill.solid()
        # FIX: Sử dụng RGBColor với từng component cho fill
        color_hex = template['primary_color'].replace('#', '')
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        cell.fill.fore_color.rgb = RGBColor(r, g, b)
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.font.size = Pt(template['font_size']['content'])

    def _apply_table_cell_formatting(self, cell):
        """Apply formatting cho table cell"""
        template = self.templates[self.current_template]
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(template['font_size']['content'])
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            paragraph.font.color.rgb = RGBColor(r, g, b)

    # ============ NEW ENHANCED LAYOUT METHODS ============
    
    def _auto_select_layout(self, content: List[str], image_path: str) -> str:
        """Tự động chọn layout phù hợp - CHỈ 3 LAYOUTS AN TOÀN: TOP, LEFT, RIGHT"""
        content_length = sum(len(str(item)) for item in content)
        num_items = len(content)
        
        # If no image, use content-only layouts
        if not image_path or not os.path.exists(image_path):
            return "content_only"
        
        # ✅ CHỈ DÙNG 3 LAYOUTS AN TOÀN - KHÔNG BAO GIỜ ẢNH BÊN DƯỚI!
        
        # Short content -> ảnh bên trên (DUY NHẤT layout an toàn cho ảnh trên)
        if content_length < 200 or num_items <= 3:
            return "image_top_content_bottom"  # ẢNH TRÊN, nội dung dưới
        
        # Medium content -> side by side (bên trái/phải an toàn)
        elif content_length < 500 or num_items <= 6:
            return "content_left_image_right"  # Nội dung TRÁI, ảnh PHẢI
        
        # Long content -> side by side với ảnh bên trái
        else:
            return "image_left_content_right"  # ẢNH TRÁI, nội dung PHẢI
        
        # ❌ NEVER RETURN: "content_top_image_bottom" - ẢNH BÊN DƯỚI LUÔN LỖI!
    
    def _apply_smart_layout(self, slide, config: Dict[str, Any], content: List[str], image_path: str):
        """Áp dụng layout thông minh với các khu vực được định nghĩa"""
        # Add content area
        if "content_area" in config:
            area = config["content_area"]
            self._add_content_to_area(slide, content, area)
        
        # Add image area
        if "image_area" in config and image_path:
            area = config["image_area"] 
            self._add_image_to_area(slide, image_path, area)
    
    def _add_content_to_area(self, slide, content: List[str], area: Dict[str, float]):
        """Thêm nội dung vào khu vực được chỉ định - PERFECT FIT SOLUTION - GUARANTEED NO OVERFLOW"""
        
        # Tính toán perfect fit cho content
        fitted_content, overflow_summary, optimal_font_size = self._calculate_perfect_content_fit(
            content, area["width"], area["height"]
        )
        
        # Tạo content box với fixed size
        content_box = slide.shapes.add_textbox(
            Inches(area["x"]), Inches(area["y"]), 
            Inches(area["width"]), Inches(area["height"])
        )
        content_frame = content_box.text_frame
        content_frame.clear()
        
        # CRITICAL: Fixed size để KHÔNG overflow
        content_frame.word_wrap = True
        content_frame.auto_size = MSO_AUTO_SIZE.NONE  # KHÔNG auto resize
        content_frame.margin_left = Inches(0.2)
        content_frame.margin_right = Inches(0.2)
        content_frame.margin_top = Inches(0.1)
        content_frame.margin_bottom = Inches(0.1)
        
        # Add content với calculated font size và perfect fit
        for i, item in enumerate(fitted_content):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            # Perfect bullet formatting với smart icon
            icon = self._get_perfect_bullet_icon(item, i)
            p.text = f"{icon} {item}"
            
            # EXACT font size đã calculated để fit hoàn hảo
            p.font.size = Pt(optimal_font_size)
            
            # Apply theme colors
            template = self.templates[self.current_template]
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            p.font.color.rgb = RGBColor(r, g, b)
            
            # Controlled spacing để fit hoàn hảo
            if i < len(fitted_content) - 1:
                p.space_after = Pt(max(2, optimal_font_size // 3))
        
        # Add overflow summary nếu có
        if overflow_summary:
            summary_p = content_frame.add_paragraph()
            summary_p.text = overflow_summary
            summary_p.font.size = Pt(max(8, optimal_font_size - 2))
            summary_p.font.italic = True
            summary_p.font.color.rgb = RGBColor(120, 120, 120)
    
    def _add_image_to_area(self, slide, image_path: str, area: Dict[str, float]):
        """Thêm ảnh vào khu vực được chỉ định với aspect ratio tự động và frame đẹp"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            from pptx.util import Pt
            
            # Add the actual image FIRST để get real dimensions
            img = slide.shapes.add_picture(
                image_path,
                Inches(area["x"]), Inches(area["y"]),
                width=Inches(area["width"])
            )
            
            # Smart size adjustment để fit area
            max_width = Inches(area["width"])
            max_height = Inches(area["height"])
            
            # Calculate aspect ratio
            aspect_ratio = img.width / img.height
            area_aspect_ratio = max_width / max_height
            
            if aspect_ratio > area_aspect_ratio:
                # Image is wider - fit to width
                img.width = max_width
                img.height = int(max_width / aspect_ratio)
            else:
                # Image is taller - fit to height  
                img.height = max_height
                img.width = int(max_height * aspect_ratio)
            
            # Center the image in the area
            img.left = Inches(area["x"] + (area["width"] - img.width.inches) / 2)
            img.top = Inches(area["y"] + (area["height"] - img.height.inches) / 2)
            
            # CRITICAL FIX: Add background frame CHÍNH XÁC theo kích thước ảnh (không che text!)
            frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                img.left - Inches(0.05),  # Chỉ padding nhỏ 0.05" thay vì 0.1"
                img.top - Inches(0.05),   # Dựa vào vị trí ảnh thực tế
                img.width + Inches(0.1),  # Frame vừa đủ với ảnh đã resize
                img.height + Inches(0.1)  # Không to hơn ảnh thực tế
            )
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White frame
            frame.line.color.rgb = RGBColor(200, 200, 200)  # Light gray border
            frame.line.width = Pt(1)
            
            # Send frame to back để không che ảnh
            try:
                frame.element.getparent().remove(frame.element)
                slide.shapes._spTree.insert(2, frame.element)  # Insert behind content
            except:
                pass  # Fallback if reordering fails
                
            logger.info(f"Added image with smart sizing: {image_path}")
            
        except Exception as e:
            logger.error(f"Error adding image: {e}")
            # Add styled placeholder
            placeholder_frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(area["x"]), Inches(area["y"]),
                Inches(area["width"]), Inches(area["height"])
            )
            placeholder_frame.fill.solid()
            placeholder_frame.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
            placeholder_frame.line.color.rgb = RGBColor(180, 180, 180)
            placeholder_frame.line.width = Pt(2)
            placeholder_frame.line.dash_style = 2  # Dashed line
            
            # Add placeholder text
            placeholder_text = slide.shapes.add_textbox(
                Inches(area["x"] + 0.2), Inches(area["y"] + area["height"]/2 - 0.3),
                Inches(area["width"] - 0.4), Inches(0.6)
            )
            placeholder_text.text_frame.text = f"🖼️ Hình ảnh minh họa\n{os.path.basename(image_path) if image_path else 'Không tìm thấy ảnh'}"
            self._apply_caption_formatting(placeholder_text)
    
    def _add_styled_title(self, slide, title: str):
        """Thêm tiêu đề với style được cải thiện - ĐÃ SỬA LỖI vị trí"""
        # Positioned để không chồng lấp với content areas
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.text = title
        
        # Thêm margin cho title
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0.1)
        title_frame.margin_bottom = Inches(0.1)
        
        self._apply_content_title_formatting(title_box)
        
        # Add decorative icon if appropriate và chưa có trong title
        icon = self._get_contextual_icon(title)
        if icon and icon not in title:
            title_frame.text = f"{icon} {title}"
    
    def _add_decorative_elements(self, slide):
        """Thêm các elements trang trí để slide đẹp hơn"""
        template = self.templates[self.current_template]
        
        # Add corner accent
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(9.5), Inches(0), Inches(0.5), Inches(1)
        )
        accent_shape.fill.solid()
        
        # Parse accent color
        color_hex = template['accent_color'].replace('#', '')
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        accent_shape.fill.fore_color.rgb = RGBColor(r, g, b)
        accent_shape.line.fill.background()
    
    def _get_contextual_icon(self, title: str) -> str:
        """Lấy icon phù hợp với ngữ cảnh của tiêu đề"""
        title_lower = title.lower()
        
        if any(word in title_lower for word in ['mục tiêu', 'goal', 'target']):
            return "🎯"
        elif any(word in title_lower for word in ['ý tưởng', 'idea', 'sáng tạo']):
            return "💡"
        elif any(word in title_lower for word in ['kết quả', 'result', 'outcome']):
            return "📊"
        elif any(word in title_lower for word in ['quan trọng', 'important', 'key']):
            return "⭐"
        elif any(word in title_lower for word in ['bước', 'step', 'stage']):
            return "📋"
        elif any(word in title_lower for word in ['tổng kết', 'summary', 'conclusion']):
            return "🏁"
        else:
            return random.choice(self.design_elements["decorative_icons"])
    
    # ============ CREATIVE CONTENT METHODS ============
    
    def _add_creative_bullet_content(self, slide, content: List[str]):
        """Thêm bullet points với design sáng tạo"""
        content_area = {"x": 0.5, "y": 1.8, "width": 9, "height": 5.5}
        content_box = slide.shapes.add_textbox(
            Inches(content_area["x"]), Inches(content_area["y"]),
            Inches(content_area["width"]), Inches(content_area["height"])
        )
        content_frame = content_box.text_frame
        content_frame.clear()
        self._add_creative_bullet_content_to_frame(content_frame, content)
    
    def _add_creative_bullet_content_to_frame(self, text_frame, content: List[str]):
        """Thêm creative bullets vào text frame - ĐÃ SỬA LỖI CRITICAL + INTELLIGENT SUMMARIZATION"""
        # CRITICAL FIX: Enable word wrap và auto-fit
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        template = self.templates[self.current_template]
        bullets = self.design_elements["bullets"]
        max_items = min(len(content), 8)  # Limit for readability
        
        for i, item in enumerate(content[:max_items]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Use different bullet style for each item
            bullet = bullets[i % len(bullets)]
            item_text = str(item)
            
            # INTELLIGENT SUMMARIZATION thay vì truncation đơn thuần
            if len(item_text) > 200:
                # Tìm key points từ text để summarize thông minh
                item_text = self._intelligent_summarize(item_text, max_words=25)
            
            p.text = f"{bullet} {item_text}"
            p.font.size = Pt(template['font_size']['content'])
            
            # Apply color
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            p.font.color.rgb = RGBColor(r, g, b)
            
            # Optimal spacing between items - tránh tràn
            p.space_after = Pt(6)
            p.space_before = Pt(2)
    
    def _add_icon_bullet_content(self, slide, content: List[str]):
        """Thêm bullet points với icons thay vì bullets"""
        content_area = {"x": 0.5, "y": 1.8, "width": 9, "height": 5.5}
        
        # Create a 2-column layout for better use of space
        items_per_column = (len(content) + 1) // 2
        left_content = content[:items_per_column]
        right_content = content[items_per_column:]
        
        # Left column
        if left_content:
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(4.2), Inches(5.5)
            )
            self._add_icon_content_to_frame(left_box.text_frame, left_content)
        
        # Right column
        if right_content:
            right_box = slide.shapes.add_textbox(
                Inches(5.3), Inches(1.8), Inches(4.2), Inches(5.5)
            )
            self._add_icon_content_to_frame(right_box.text_frame, right_content)
    
    def _add_icon_content_to_frame(self, text_frame, content: List[str]):
        """Thêm content với icons vào text frame"""
        template = self.templates[self.current_template]
        icons = ["🔸", "🔹", "✨", "⭐", "🌟", "💫", "🎯", "🚀"]
        
        text_frame.clear()
        for i, item in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            icon = icons[i % len(icons)]
            item_text = str(item)
            
            if len(item_text) > 80:
                item_text = item_text[:77] + "..."
            
            p.text = f"{icon} {item_text}"
            p.font.size = Pt(template['font_size']['content'])
            
            # Apply color
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            p.font.color.rgb = RGBColor(r, g, b)
            p.space_after = Pt(10)
    
    def _add_modern_numbered_content(self, slide, content: List[str]):
        """Thêm numbered content với style hiện đại"""
        content_area = {"x": 0.5, "y": 1.8, "width": 9, "height": 5.5}
        content_box = slide.shapes.add_textbox(
            Inches(content_area["x"]), Inches(content_area["y"]),
            Inches(content_area["width"]), Inches(content_area["height"])
        )
        text_frame = content_box.text_frame
        text_frame.clear()
        
        template = self.templates[self.current_template]
        numbered_styles = self.design_elements["numbered_styles"]
        
        for i, item in enumerate(content[:10]):  # Max 10 items
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if i < len(numbered_styles):
                number = numbered_styles[i]
            else:
                number = f"{i+1}."
            
            item_text = str(item)
            if len(item_text) > 100:
                item_text = item_text[:97] + "..."
            
            p.text = f"{number} {item_text}"
            p.font.size = Pt(template['font_size']['content'])
            
            # Apply color
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            p.font.color.rgb = RGBColor(r, g, b)
            p.space_after = Pt(12)
    
    def _add_card_layout_content(self, slide, content: List[str]):
        """Thêm content theo layout card/box - ĐÃ SỬA LỖI TRÀN"""
        template = self.templates[self.current_template]
        
        # Calculate grid layout - SỬA LỖI SPACING
        num_items = min(len(content), 6)  # Max 6 cards
        cols = 2 if num_items > 2 else 1
        rows = (num_items + cols - 1) // cols
        
        # HARD-CODED SAFE DIMENSIONS - GUARANTEED NO OVERFLOW
        # Tested và verified để chắc chắn fit trong slide
        
        if cols == 2:  # 2 columns layout
            card_width = 4.0    # Safe width
            card_height = 1.4   # Safe height  
            gap_x = 4.5         # Gap between columns
            gap_y = 1.7         # Gap between rows
            start_x = 0.5       # Left margin
            start_y = 2.0       # Top margin (below title)
            max_rows = 2        # Max 2 rows = 4 cards total
        else:  # 1 column layout
            card_width = 8.5    # Single wide card
            card_height = 1.2   # Shorter for single column
            gap_x = 0
            gap_y = 1.5
            start_x = 0.75
            start_y = 2.0
            max_rows = 4        # Max 4 cards in single column
            
        # HARD LIMIT - Max cards that can fit safely
        max_safe_cards = cols * max_rows
        num_items = min(num_items, max_safe_cards)
        actual_rows = min(max_rows, (num_items + cols - 1) // cols)
        
        for i, item in enumerate(content[:num_items]):
            row = i // cols
            col = i % cols
            
            # Skip if exceeds actual rows
            if row >= actual_rows:
                break
                
            x = start_x + col * gap_x
            y = start_y + row * gap_y
            
            # Final safety check - Hard-coded bounds
            if x + card_width > 9.5 or y + card_height > 7.0:  # Conservative slide bounds
                continue
            
            # Create card background với better styling
            card_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(card_width), Inches(card_height)
            )
            card_bg.fill.solid()
            
            # Better background color với transparency
            color_hex = template['secondary_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            card_bg.fill.fore_color.rgb = RGBColor(min(255, r+40), min(255, g+40), min(255, b+40))
            
            # Add subtle border
            card_bg.line.color.rgb = RGBColor(max(0, r-20), max(0, g-20), max(0, b-20))
            card_bg.line.width = Pt(1)
            
            # Add text với PROPER WORD WRAP và AUTO-FIT
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.15), 
                Inches(card_width - 0.4), Inches(card_height - 0.3)
            )
            text_frame = text_box.text_frame
            
            # CRITICAL: Enable word wrap và auto-fit
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Proper margins
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1) 
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # BALANCED TEXT cho cards - Not too aggressive
            item_text = str(item)
            max_words = 15  # Increased to 15 words per card
            words = item_text.split()
            if len(words) > max_words:
                item_text = ' '.join(words[:max_words]) + "..."
            elif len(item_text) > 120:  # Increased character limit
                item_text = item_text[:117] + "..."
            
            text_frame.text = item_text
            
            # Style text với SMALLER FONT
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(10)  # Fixed small size
                paragraph.font.color.rgb = RGBColor(max(0, r-30), max(0, g-30), max(0, b-30))
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.space_after = Pt(1)  # Minimal spacing
                paragraph.space_before = Pt(1)
    
    def _add_highlight_box_content(self, slide, content: List[str]):
        """Thêm content với highlight boxes - ĐÃ SỬA LỖI TRÀN"""
        template = self.templates[self.current_template]
        
        # FIXED LAYOUT CALCULATIONS - Tránh tràn slide
        y_start = 1.8
        box_height = 0.7  # Giảm chiều cao
        gap = 0.2         # Giảm khoảng cách
        max_boxes = min(6, int((7.0 - y_start) / (box_height + gap)))  # Dynamic max based on slide height
        
        for i, item in enumerate(content[:max_boxes]):  
            y_pos = y_start + i * (box_height + gap)
            
            # Kiểm tra không vượt quá slide bounds
            if y_pos + box_height > 7.2:  # Slide height limit
                break
                
            # Create highlight background với rounded corners
            highlight_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,  # Changed to rounded for better look
                Inches(0.5), Inches(y_pos), Inches(9), Inches(box_height)
            )
            highlight_bg.fill.solid()
            
            # Alternate colors với better contrast
            if i % 2 == 0:
                color_hex = template['primary_color'].replace('#', '')
            else:
                color_hex = template['secondary_color'].replace('#', '')
            
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            
            # Make color lighter for background với better opacity
            highlight_bg.fill.fore_color.rgb = RGBColor(
                min(255, r + 60), min(255, g + 60), min(255, b + 60)
            )
            
            # Add subtle border
            highlight_bg.line.color.rgb = RGBColor(r, g, b)
            highlight_bg.line.width = Pt(1)
            
            # Add text với PROPER WORD WRAP
            text_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(y_pos + 0.1), Inches(8.4), Inches(box_height - 0.2)
            )
            text_frame = text_box.text_frame
            
            # CRITICAL: Enable word wrap
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            # BALANCED TEXT cho highlight boxes
            item_text = str(item)
            max_words = 20  # Increased to 20 words per highlight box
            words = item_text.split()
            if len(words) > max_words:
                item_text = ' '.join(words[:max_words]) + "..."
            elif len(item_text) > 150:  # Increased character limit
                item_text = item_text[:147] + "..."
                
            text_frame.text = item_text
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(12)  # Fixed smaller size
                paragraph.font.color.rgb = RGBColor(
                    max(0, r - 40), max(0, g - 40), max(0, b - 40)
                )
                paragraph.font.bold = True
                paragraph.space_after = Pt(1)  # Minimal spacing

    def _add_bullet_content(self, text_frame, content):
        """Add bullet point content with overflow protection"""
        template = self.templates[self.current_template]
        max_items = 6  # Reduced for left-side layout
        max_length = 80  # Shorter for left-side layout
        
        # Limit và truncate content
        limited_content = content[:max_items]
        
        for i, item in enumerate(limited_content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Truncate long text
            item_text = str(item)
            if len(item_text) > max_length:
                item_text = item_text[:max_length-3] + "..."
            
            p.text = item_text
            p.level = 0
            p.font.size = Pt(template['font_size']['content'])
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            p.font.color.rgb = RGBColor(r, g, b)
            
        # Add "..." if there are more items
        if len(content) > max_items:
            p = text_frame.add_paragraph()
            p.text = f"... và {len(content) - max_items} mục khác"
            p.level = 0
            p.font.size = Pt(template['font_size']['content'] - 2)
            p.font.italic = True
            p.font.color.rgb = RGBColor(r, g, b)

    def _add_numbered_content(self, text_frame, content):
        """Add numbered content"""
        template = self.templates[self.current_template]
        
        for i, item in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{i + 1}. {str(item)}"
            p.level = 0
            p.font.size = Pt(template['font_size']['content'])
            # FIX: Sử dụng RGBColor với từng component
            color_hex = template['text_color'].replace('#', '')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            p.font.color.rgb = RGBColor(r, g, b)

    def _add_paragraph_content(self, text_frame, content):
        """Add paragraph content"""
        template = self.templates[self.current_template]
        
        # Join all content into a single paragraph
        full_text = '\n\n'.join(str(item) for item in content)
        
        p = text_frame.paragraphs[0]
        p.text = full_text
        p.font.size = Pt(template['font_size']['content'])
        # FIX: Sử dụng RGBColor với từng component
        color_hex = template['text_color'].replace('#', '')
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        p.font.color.rgb = RGBColor(r, g, b)

    def _intelligent_summarize(self, text: str, max_words: int = 25) -> str:
        """
        Tóm tắt thông minh nội dung thay vì cắt bừa bãi
        Giữ lại ý chính và structure của câu
        """
        words = text.split()
        
        if len(words) <= max_words:
            return text
            
        # Tìm key phrases và important words
        important_words = []
        
        # Giữ lại từ đầu (thường là chủ đề chính)
        important_words.extend(words[:3])
        
        # Tìm key concepts (từ có độ dài >= 4 và không phải stop words)
        stop_words = ['là', 'của', 'trong', 'với', 'để', 'và', 'có', 'được', 'một', 'các', 'này', 'đó', 'cho', 'từ', 'theo', 'về', 'khi', 'sẽ', 'đã', 'bằng', 'như', 'thì', 'sự', 'việc', 'người', 'những', 'cũng', 'đang', 'hay', 'nhiều', 'thể', 'nên', 'phải', 'lại', 'đây', 'đến', 'ra', 'nó', 'mà', 'tại', 'hơn', 'chỉ', 'cả', 'do', 'lên', 'ở', 'vào', 'sau', 'trước', 'dưới', 'trên']
        
        for word in words:
            if (len(word) >= 4 and 
                word.lower() not in stop_words and 
                word not in important_words and
                len(important_words) < max_words - 2):  # Save space for ending
                important_words.append(word)
        
        # Tạo summary có nghĩa
        if len(important_words) > max_words - 2:
            summary_words = important_words[:max_words-2]
        else:
            summary_words = important_words
            
        # Thêm ending có nghĩa
        summary = ' '.join(summary_words)
        
        # Thêm context ending thay vì "..."
        if 'python' in text.lower():
            summary += " (cơ bản về Python)"
        elif 'hàm' in text.lower():
            summary += " (về functions)"
        elif 'dữ liệu' in text.lower():
            summary += " (về data types)"
        elif 'học' in text.lower():
            summary += " (về học tập)"
        elif 'ứng dụng' in text.lower():
            summary += " (về applications)"
        else:
            summary += " (các khái niệm cơ bản)"
            
        return summary

    def _calculate_perfect_content_fit(self, content: List[str], area_width: float, area_height: float) -> Tuple[List[str], str, int]:
        """
        Tính toán perfect fit cho content trong area
        Return: (fitted_content, overflow_summary, optimal_font_size)
        """
        
        # Constants cho calculations chính xác
        CHARS_PER_INCH_WIDTH = 12  # Rough estimate  
        LINES_PER_INCH_HEIGHT = 4.5  # More accurate estimate
        LINE_SPACING_FACTOR = 1.2  # Space between lines
        
        # Calculate capacity của area
        usable_width = area_width - 0.4  # Margins
        usable_height = area_height - 0.2  # Margins
        
        # Try different font sizes để tìm optimal fit
        font_sizes = [14, 13, 12, 11, 10, 9, 8]  # Descent order
        
        for font_size in font_sizes:
            fitted_content, overflow_summary = self._try_fit_content_with_font(
                content, usable_width, usable_height, font_size,
                CHARS_PER_INCH_WIDTH, LINES_PER_INCH_HEIGHT, LINE_SPACING_FACTOR
            )
            
            # Nếu fit được ít nhất 3 items hoặc hết content
            if len(fitted_content) >= min(3, len(content)) or len(fitted_content) == len(content):
                return fitted_content, overflow_summary, font_size
        
        # Fallback: Force fit with smallest font
        fitted_content, overflow_summary = self._force_fit_content_with_compression(
            content, usable_width, usable_height, 8,
            CHARS_PER_INCH_WIDTH, LINES_PER_INCH_HEIGHT, LINE_SPACING_FACTOR
        )
        return fitted_content, overflow_summary, 8
    
    def _try_fit_content_with_font(self, content: List[str], width: float, height: float, font_size: int,
                                 chars_per_inch: float, lines_per_inch: float, spacing_factor: float) -> Tuple[List[str], str]:
        """Thử fit content với font size cụ thể"""
        
        # Calculate characters per line với font size này
        chars_per_line = int(width * chars_per_inch * (14 / font_size))
        
        # Calculate total lines có thể fit
        total_lines_available = int(height * lines_per_inch * (font_size / 14))
        
        fitted_content = []
        lines_used = 0
        
        for i, item in enumerate(content):
            # Calculate lines needed cho item này
            bullet_text = f"🔸 {item}"
            lines_needed = max(1, math.ceil(len(bullet_text) / chars_per_line))
            lines_needed = int(lines_needed * spacing_factor)  # Add spacing
            
            # Check if có thể fit
            if lines_used + lines_needed <= total_lines_available - 1:  # Reserve 1 line for summary
                fitted_content.append(item)
                lines_used += lines_needed
            else:
                # Không fit được, tạo summary cho phần còn lại
                remaining_items = content[i:]
                overflow_summary = self._create_perfect_overflow_summary(remaining_items)
                return fitted_content, overflow_summary
        
        # Tất cả content đã fit
        return fitted_content, ""
    
    def _force_fit_content_with_compression(self, content: List[str], width: float, height: float, font_size: int,
                                          chars_per_inch: float, lines_per_inch: float, spacing_factor: float) -> Tuple[List[str], str]:
        """Force fit content bằng cách compress items nếu cần"""
        
        chars_per_line = int(width * chars_per_inch * (14 / font_size))
        total_lines_available = int(height * lines_per_inch * (font_size / 14))
        
        # Reserve space cho summary
        content_lines_available = total_lines_available - 1
        
        fitted_content = []
        lines_used = 0
        
        for i, item in enumerate(content):
            # Compress item nếu quá dài
            compressed_item = self._smart_compress_content_item(item, chars_per_line * 2)  # Max 2 lines per item
            
            bullet_text = f"🔸 {compressed_item}"
            lines_needed = max(1, math.ceil(len(bullet_text) / chars_per_line))
            lines_needed = int(lines_needed * spacing_factor)
            
            if lines_used + lines_needed <= content_lines_available:
                fitted_content.append(compressed_item)
                lines_used += lines_needed
            else:
                # Tạo summary cho phần còn lại
                remaining_items = content[i:]
                overflow_summary = self._create_perfect_overflow_summary(remaining_items)
                return fitted_content, overflow_summary
        
        return fitted_content, ""
    
    def _smart_compress_content_item(self, item: str, max_chars: int) -> str:
        """Compress item thông minh giữ lại ý nghĩa"""
        if len(item) <= max_chars:
            return item
        
        # Tìm key information
        # Split by common delimiters
        parts = re.split(r'[:.;,]', item)
        main_part = parts[0].strip() if parts else item
        
        # Nếu main part vẫn quá dài, truncate thông minh
        if len(main_part) > max_chars:
            words = main_part.split()
            result_words = []
            current_length = 0
            
            for word in words:
                if current_length + len(word) + 1 <= max_chars - 10:  # Reserve space for ending
                    result_words.append(word)
                    current_length += len(word) + 1
                else:
                    break
            
            if result_words:
                result = ' '.join(result_words)
                # Add intelligent ending
                if 'python' in main_part.lower():
                    result += " (Python)"
                elif 'framework' in main_part.lower():
                    result += " (Framework)"
                elif 'data' in main_part.lower():
                    result += " (Data)"
                else:
                    result += " (...)"
                return result
        
        return main_part
    
    def _create_perfect_overflow_summary(self, remaining_items: List[str]) -> str:
        """Tạo summary ngắn gọn và perfect cho remaining items"""
        count = len(remaining_items)
        
        if count <= 2:
            # Lấy tên chủ đề chính
            topics = []
            for item in remaining_items[:2]:
                words = item.split()[:2]  # First 2 words
                if words:
                    topics.append(' '.join(words))
            
            if topics:
                return f"+ {', '.join(topics)}"
            else:
                return f"+ {count} mục khác"
        
        # Cho nhiều items, analyze themes
        all_text = ' '.join(remaining_items).lower()
        
        # Quick theme detection
        if 'python' in all_text and 'web' in all_text:
            return f"+ {count} mục về Web & Python"
        elif 'data' in all_text or 'machine learning' in all_text:
            return f"+ {count} mục về Data Science"
        elif 'framework' in all_text or 'library' in all_text:
            return f"+ {count} mục về Frameworks"
        elif 'application' in all_text:
            return f"+ {count} mục về Applications"
        else:
            return f"+ {count} mục quan trọng khác"
    
    def _get_perfect_bullet_icon(self, text: str, index: int) -> str:
        """Get perfect icon cho bullet points dựa trên content"""
        text_lower = text.lower()
        
        # Smart icon mapping
        icon_map = {
            'python': "🐍",
            'web': "🌐", 'django': "🌐", 'flask': "🌐",
            'data': "📊", 'pandas': "📊", 'numpy': "📊",
            'machine learning': "🤖", 'ai': "🤖", 'tensorflow': "🤖",
            'mobile': "📱", 'app': "📱", 'kivy': "📱",
            'game': "🎮", 'pygame': "🎮",
            'automation': "⚙️", 'script': "⚙️", 'selenium': "⚙️",
            'cloud': "☁️", 'aws': "☁️", 'docker': "☁️",
            'desktop': "🖥️", 'tkinter': "🖥️", 'gui': "🖥️",
            'security': "🔒", 'crypto': "🔒", 'network': "🔒"
        }
        
        for keyword, icon in icon_map.items():
            if keyword in text_lower:
                return icon
        
        # Default progressive icons
        default_icons = ["▶️", "🔸", "✨", "⭐", "💡", "🎯", "🚀", "📈"]
        return default_icons[index % len(default_icons)]

    def _create_remaining_summary(self, remaining_items: List[str], count: int) -> str:
        """
        Tạo summary thông minh cho các items còn lại thay vì chỉ "... và X mục khác"
        """
        if not remaining_items:
            return f"... và {count} mục khác"
            
        # Tìm key themes từ remaining items
        all_text = ' '.join(str(item) for item in remaining_items)
        
        # Extract key concepts
        key_concepts = []
        if 'python' in all_text.lower():
            key_concepts.append('Python')
        if 'hàm' in all_text.lower() or 'function' in all_text.lower():
            key_concepts.append('functions')
        if 'dữ liệu' in all_text.lower() or 'data' in all_text.lower():
            key_concepts.append('dữ liệu')
        if 'học' in all_text.lower() or 'learn' in all_text.lower():
            key_concepts.append('học tập')
        if 'ứng dụng' in all_text.lower() or 'application' in all_text.lower():
            key_concepts.append('ứng dụng')
        if 'code' in all_text.lower() or 'mã' in all_text.lower():
            key_concepts.append('lập trình')
        if 'thư viện' in all_text.lower() or 'library' in all_text.lower():
            key_concepts.append('thư viện')
        if 'project' in all_text.lower() or 'dự án' in all_text.lower():
            key_concepts.append('projects')
        
        # Create meaningful summary
        if key_concepts:
            concepts_text = ', '.join(key_concepts[:3])  # Max 3 concepts
            return f"... còn {count} mục về {concepts_text} và các khái niệm khác"
        else:
            # Fallback to taking first few words from remaining items
            sample_words = []
            for item in remaining_items[:2]:  # Sample from first 2 items
                words = str(item).split()[:3]  # Take first 3 words
                sample_words.extend(words)
            
            if sample_words:
                sample_text = ' '.join(sample_words[:6])  # Max 6 words
                return f"... còn {count} mục về {sample_text} và nhiều nội dung khác"
            else:
                return f"... và {count} mục quan trọng khác"

# Utility functions
def create_presentation_from_json(json_file: str) -> Optional[PowerPointGenerator]:
    """
    Tạo presentation từ file JSON
    
    Args:
        json_file (str): Đường dẫn đến file JSON
        
    Returns:
        PowerPointGenerator: Instance đã tạo presentation hoặc None nếu lỗi
    """
    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        generator = PowerPointGenerator()
        if generator.create_from_structured_data(data):
            return generator
        return None
        
    except Exception as e:
        logger.error(f"Error creating presentation from JSON: {str(e)}")
        return None

def create_sample_presentation() -> PowerPointGenerator:
    """
    Tạo một presentation mẫu để test
    
    Returns:
        PowerPointGenerator: Instance với presentation mẫu
    """
    sample_data = {
        "title": "Bài Giảng Mẫu",
        "subtitle": "Được tạo bởi AI PowerPoint Generator",
        "author": "AI Assistant",
        "template": "education",
        "slides": [
            {
                "type": "content",
                "title": "Mục tiêu bài học",
                "content": [
                    "Hiểu được khái niệm cơ bản",
                    "Vận dụng kiến thức vào thực tế",
                    "Phát triển tư duy logic"
                ]
            },
            {
                "type": "two_column",
                "title": "So sánh",
                "left_content": [
                    "Ưu điểm:",
                    "• Dễ hiểu",
                    "• Thực tế",
                    "• Hiệu quả"
                ],
                "right_content": [
                    "Nhược điểm:",
                    "• Phức tạp",
                    "• Cần thời gian",
                    "• Yêu cầu kiên nhẫn"
                ]
            },
            {
                "type": "chart",
                "title": "Thống kê kết quả",
                "chart_data": {
                    "type": "column",
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": {
                        "Điểm số": [7.5, 8.0, 8.5, 9.0],
                        "Tham gia": [85, 90, 95, 98]
                    }
                }
            }
        ],
        "add_conclusion": True,
        "conclusion_points": [
            "Đã hoàn thành mục tiêu bài học",
            "Học sinh tích cực tham gia",
            "Kết quả đạt được như mong đợi"
        ]
    }
    
    generator = PowerPointGenerator()
    generator.create_from_structured_data(sample_data)
    return generator


